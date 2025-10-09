import os
import hashlib
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from difflib import SequenceMatcher
from PIL import Image
from io import BytesIO
from collections import defaultdict

def hash_image(image_blob):
    """Return SHA256 hash of image bytes."""
    return hashlib.sha256(image_blob).hexdigest()

def compare_text(old_text, new_text):
    """Return HTML diff of two text blocks."""
    sm = SequenceMatcher(None, old_text, new_text)
    html = ""
    for tag, i1, i2, j1, j2 in sm.get_opcodes():
        if tag == "equal":
            html += old_text[i1:i2]
        elif tag == "replace":
            html += f'<span style="background:yellow;">{new_text[j1:j2]}</span>'
        elif tag == "delete":
            html += f'<span style="background:#faa;color:red;text-decoration:line-through;">{old_text[i1:i2]}</span>'
        elif tag == "insert":
            html += f'<span style="background:#cfc;color:green;">{new_text[j1:j2]}</span>'
    return html

def extract_slide_content(slide):
    """Extract text, notes, images, and tables from a slide."""
    content = {
        "text": [],
        "notes": "",
        "images": [],
        "tables": [],
        "formatting": []
    }
    # Text and formatting
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            para_text = paragraph.text
            runs = []
            for run in paragraph.runs:
                runs.append({
                    "text": run.text,
                    "bold": run.font.bold,
                    "italic": run.font.italic,
                    "color": run.font.color.rgb if run.font.color and run.font.color.rgb else None,
                    "font": run.font.name
                })
            content["text"].append(para_text)
            content["formatting"].append(runs)
    # Notes
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        content["notes"] = slide.notes_slide.notes_text_frame.text
    # Images
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_blob = shape.image.blob
            image_hash = hash_image(image_blob)
            content["images"].append({
                "hash": image_hash,
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            })
    # Tables
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_data = []
            table = shape.table
            for row in table.rows:
                table_data.append([cell.text for cell in row.cells])
            content["tables"].append(table_data)
    return content

def compare_images(old_imgs, new_imgs):
    """Compare image lists by hash and position."""
    old_set = {(img['hash'], img['left'], img['top'], img['width'], img['height']) for img in old_imgs}
    new_set = {(img['hash'], img['left'], img['top'], img['width'], img['height']) for img in new_imgs}
    added = new_set - old_set
    removed = old_set - new_set
    unchanged = old_set & new_set
    return added, removed, unchanged

def compare_tables(old_tables, new_tables):
    """Compare tables by structure and content."""
    diffs = []
    max_len = max(len(old_tables), len(new_tables))
    for i in range(max_len):
        old = old_tables[i] if i < len(old_tables) else []
        new = new_tables[i] if i < len(new_tables) else []
        if old != new:
            diffs.append((i, old, new))
    return diffs

def compare_formatting(old_fmt, new_fmt):
    """Compare formatting runs (bold, italic, color, font)."""
    diffs = []
    for i, (old_runs, new_runs) in enumerate(zip(old_fmt, new_fmt)):
        if old_runs != new_runs:
            diffs.append((i, old_runs, new_runs))
    return diffs

def compare_slides(old_slide, new_slide):
    """Compare all aspects of two slides."""
    result = {}
    # Text
    old_text = "\n".join(old_slide["text"])
    new_text = "\n".join(new_slide["text"])
    result["text_diff"] = compare_text(old_text, new_text)
    # Notes
    result["notes_diff"] = compare_text(old_slide["notes"], new_slide["notes"])
    # Images
    added, removed, unchanged = compare_images(old_slide["images"], new_slide["images"])
    result["images_added"] = list(added)
    result["images_removed"] = list(removed)
    # Tables
    result["table_diffs"] = compare_tables(old_slide["tables"], new_slide["tables"])
    # Formatting
    result["formatting_diffs"] = compare_formatting(old_slide["formatting"], new_slide["formatting"])
    return result

def generate_html_report(slide_diffs, summary, old_ppt, new_ppt):
    html = [
        "<html><head><meta charset='utf-8'><title>PPT Comparison Report</title>",
        "<style>",
        "body{font-family:Arial,sans-serif;}",
        ".slide-section{margin-bottom:40px;}",
        ".diff-block{border:1px solid #ccc;padding:10px;margin:10px 0;}",
        ".added{background:#cfc;}",
        ".removed{background:#faa;text-decoration:line-through;}",
        ".changed{background:yellow;}",
        ".summary-table{border-collapse:collapse; margin-bottom:20px;}",
        ".summary-table td, .summary-table th{border:1px solid #ccc;padding:4px 8px;}",
        "</style></head><body>",
        f"<h1>PPT Comparison Report</h1>",
        f"<h2>Old File: {old_ppt}</h2>",
        f"<h2>New File: {new_ppt}</h2>",
        "<h2>Summary of Changes</h2>",
        summary,
    ]
    for idx, diff in enumerate(slide_diffs):
        html.append(f"<div class='slide-section'><h2>Slide {idx+1}</h2>")
        html.append("<div class='diff-block'><b>Text Diff:</b><br>" + diff["text_diff"] + "</div>")
        html.append("<div class='diff-block'><b>Notes Diff:</b><br>" + diff["notes_diff"] + "</div>")
        # Images
        html.append("<div class='diff-block'><b>Images Added:</b><br>")
        for img in diff["images_added"]:
            html.append(f"<span class='added'>Image hash: {img[0]}, Pos: {img[1:]}</span><br>")
        html.append("</div>")
        html.append("<div class='diff-block'><b>Images Removed:</b><br>")
        for img in diff["images_removed"]:
            html.append(f"<span class='removed'>Image hash: {img[0]}, Pos: {img[1:]}</span><br>")
        html.append("</div>")
        # Tables
        if diff["table_diffs"]:
            html.append("<div class='diff-block'><b>Table Differences:</b><br>")
            for i, old, new in diff["table_diffs"]:
                html.append(f"<span class='removed'>Old Table {i+1}: {old}</span><br>")
                html.append(f"<span class='added'>New Table {i+1}: {new}</span><br>")
            html.append("</div>")
        # Formatting
        if diff["formatting_diffs"]:
            html.append("<div class='diff-block'><b>Formatting Differences:</b><br>")
            for i, old, new in diff["formatting_diffs"]:
                html.append(f"<span class='removed'>Old Paragraph {i+1}: {old}</span><br>")
                html.append(f"<span class='added'>New Paragraph {i+1}: {new}</span><br>")
            html.append("</div>")
        html.append("</div>")
    html.append("</body></html>")
    return "\n".join(html)

def summarize_changes(slide_diffs):
    total_slides = len(slide_diffs)
    text_changes = sum(1 for d in slide_diffs if d["text_diff"].find("span") != -1)
    notes_changes = sum(1 for d in slide_diffs if d["notes_diff"].find("span") != -1)
    images_added = sum(len(d["images_added"]) for d in slide_diffs)
    images_removed = sum(len(d["images_removed"]) for d in slide_diffs)
    table_changes = sum(len(d["table_diffs"]) for d in slide_diffs)
    formatting_changes = sum(len(d["formatting_diffs"]) for d in slide_diffs)
    summary = f"""
    <table class='summary-table'>
    <tr><th>Total Slides</th><td>{total_slides}</td></tr>
    <tr><th>Slides with Text Changes</th><td>{text_changes}</td></tr>
    <tr><th>Slides with Notes Changes</th><td>{notes_changes}</td></tr>
    <tr><th>Total Images Added</th><td>{images_added}</td></tr>
    <tr><th>Total Images Removed</th><td>{images_removed}</td></tr>
    <tr><th>Slides with Table Changes</th><td>{table_changes}</td></tr>
    <tr><th>Slides with Formatting Changes</th><td>{formatting_changes}</td></tr>
    </table>
    """
    return summary

def main(old_pptx, new_pptx, output_html):
    prs_old = Presentation(old_pptx)
    prs_new = Presentation(new_pptx)
    slides_old = [extract_slide_content(s) for s in prs_old.slides]
    slides_new = [extract_slide_content(s) for s in prs_new.slides]
    max_slides = max(len(slides_old), len(slides_new))
    slide_diffs = []
    for i in range(max_slides):
        old = slides_old[i] if i < len(slides_old) else {"text":[],"notes":"","images":[],"tables":[],"formatting":[]}
        new = slides_new[i] if i < len(slides_new) else {"text":[],"notes":"","images":[],"tables":[],"formatting":[]}
        slide_diffs.append(compare_slides(old, new))
    summary = summarize_changes(slide_diffs)
    html = generate_html_report(slide_diffs, summary, old_pptx, new_pptx)
    with open(output_html, "w", encoding="utf-8") as f:
        f.write(html)
    print(f"Comparison report generated: {output_html}")

if __name__ == "__main__":
    # Update these paths for your files
    old_pptx = r"path/to/your/old/presentation.pptx"
    new_pptx = r"path/to/your/new/presentation.pptx"
    output_html = r"ppt_comparison_report.html"
    main(old_pptx, new_pptx, output_html)