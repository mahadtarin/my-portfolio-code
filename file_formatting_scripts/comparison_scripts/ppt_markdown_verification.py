import os
import sys
import re
import pandas as pd
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from difflib import SequenceMatcher

# Set stdout encoding to utf-8
sys.stdout.reconfigure(encoding='utf-8')

def extract_slide_content(slide, slide_number):
    """Extracts text content (including formatting) from a PowerPoint slide."""
    slide_content = []

    # Extract slide title
    slide_title = slide.shapes.title.text if slide.shapes.title else f"Slide_{slide_number}"
    slide_content.append(f"# {slide_title}\n")

    # Extract shapes and process each one
    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape == slide.shapes.title:
                continue  # Skip the title since it's already added
            for paragraph in shape.text_frame.paragraphs:
                formatted_text = format_paragraph(paragraph)
                slide_content.append(formatted_text)

    # Extract notes (if available)
    if slide.has_notes_slide:
        slide_notes = process_slide_notes(slide)
        if slide_notes:
            slide_content.append(slide_notes)

    return "\n".join(slide_content)


def format_paragraph(paragraph):
    """Formats a paragraph's text content."""
    formatted_text = []

    for run in paragraph.runs:
        text = run.text.strip()
        if not text:
            continue

        # Apply bold and italic formatting
        if run.font.bold and run.font.italic:
            text = f"***{text}***"
        elif run.font.bold:
            text = f"**{text}**"
        elif run.font.italic:
            text = f"*{text}*"

        formatted_text.append(text)

    return " ".join(formatted_text)


def process_slide_notes(slide):
    """Extracts and formats notes from a slide."""
    if slide.notes_slide and slide.notes_slide.notes_text_frame:
        notes_frame = slide.notes_slide.notes_text_frame
        paragraphs = notes_frame.paragraphs
        formatted_notes = []

        for paragraph in paragraphs:
            buffer = []

            for run in paragraph.runs:
                text = run.text.strip()
                if not text:
                    continue

                # Apply bold and italic formatting to notes
                if run.font.bold and run.font.italic:
                    text = f"***{text}***"
                elif run.font.bold:
                    text = f"**{text}**"
                elif run.font.italic:
                    text = f"*{text}*"

                buffer.append(text)

            formatted_notes.append(" ".join(buffer))

        return f"### Notes:\n" + "\n".join(formatted_notes)

    return ""


def extract_text_similarity(text1, text2):
    print(f"Text1: {text1}")
    print(f"Text2: {text2}")
    vectorizer = TfidfVectorizer().fit_transform([text1, text2])
    vectors = vectorizer.toarray()
    cos_sim = cosine_similarity(vectors)
    return cos_sim[0][1]


def verify_format_consistency(ppt_content, md_content):
    ppt_lines = ppt_content.split("\n")
    md_lines = md_content.split("\n")

    errors = []

    # Check bullet points (lines starting with '*')
    ppt_bullets = [line.strip() for line in ppt_lines if line.strip().startswith('*')]
    md_bullets = [line.strip() for line in md_lines if line.strip().startswith('*')]

    # Soft assert bullet points count and content
    if len(ppt_bullets) != len(md_bullets):
        errors.append("Number of bullet points do not match between PPT and Markdown.")
    for ppt_bullet, md_bullet in zip(ppt_bullets, md_bullets):
        if ppt_bullet != md_bullet:
            errors.append(f"Bullet point mismatch: {ppt_bullet} != {md_bullet}")

    if errors:
        print("Formatting inconsistencies found:")
        for error in errors:
            print(f"- {error}")
        return False  # Indicate that there were errors
    else:
        print("Format consistency verified.")

def generate_diff_html(ppt_text, md_text):
    """
    Generates HTML highlighting differences between PowerPoint and Markdown content at the word level.
    Args:
        ppt_text (str): Text content from the PowerPoint slide.
        md_text (str): Text content from the Markdown file.
    Returns:
        tuple: Two HTML strings with differences highlighted.
    """
    # Split the content into words for word-level comparison
    ppt_words = ppt_text.split()
    md_words = md_text.split()

    # Use SequenceMatcher to compare word lists
    diff = SequenceMatcher(None, ppt_words, md_words).get_opcodes()
    ppt_html = ""
    md_html = ""

    for tag, i1, i2, j1, j2 in diff:
        if tag == "replace":
            ppt_html += f'<span class="removed">{" ".join(ppt_words[i1:i2])}</span> '
            md_html += f'<span class="added">{" ".join(md_words[j1:j2])}</span> '
        elif tag == "delete":
            ppt_html += f'<span class="removed">{" ".join(ppt_words[i1:i2])}</span> '
        elif tag == "insert":
            md_html += f'<span class="added">{" ".join(md_words[j1:j2])}</span> '
        elif tag == "equal":
            ppt_html += f'<span class="unchanged">{" ".join(ppt_words[i1:i2])}</span> '
            md_html += f'<span class="unchanged">{" ".join(md_words[j1:j2])}</span> '

    return ppt_html.strip(), md_html.strip()

def compare_ppt_to_markdown_advanced(ppt_content_array, md_content_dict, threshold=0.75, html_filename="similarity_results.html"):
    """
    Compares PowerPoint content to Markdown content and generates an HTML report.
    Args:
        ppt_content_array (list): List of dictionaries containing slide content.
        md_content_dict (dict): Dictionary containing Markdown content by slide number.
        threshold (float): Minimum similarity score for a confident match.
        html_filename (str): Path to the HTML report file.
    """
    # Check if ppt_content_array is a list
    if not isinstance(ppt_content_array, list):
        raise ValueError("ppt_content_array must be a list of dictionaries with slide content.")

    # Check if md_content_dict is a dictionary
    if not isinstance(md_content_dict, dict):
        raise ValueError("md_content_dict must be a dictionary with markdown content.")

    all_slides_match = True
    similarity_below_threshold = []

    # Open the HTML file for writing
    with open(html_filename, "w", encoding="utf-8") as html:
        # Write the HTML header
        html.write("""
        <!DOCTYPE html>
        <html lang="en">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>PowerPoint to Markdown Comparison Report</title>
            <style>
                body { font-family: Arial, sans-serif; line-height: 1.6; margin: 20px; }
                h1 { color: #333; }
                h2 { color: #555; }
                .low-similarity { color: red; font-weight: bold; }
                .high-similarity { color: green; font-weight: bold; }
                .content-section { margin-bottom: 20px; }
                .diff { margin: 10px 0; padding: 10px; border: 1px solid #ddd; background: #f4f4f4; }
                .ppt-content { color: blue; }
                .md-content { color: purple; }
            </style>
        </head>
        <body>
            <h1>PowerPoint to Markdown Comparison Report</h1>
        """)

        # Iterate over each slide in ppt_content_array
        for i, slide_content in enumerate(ppt_content_array):
            slide_number = i + 1

            # Combine title with the rest of the slide content
            ppt_content = slide_content['content']  # Assuming slide_content is a list of strings

            # Get corresponding Markdown content for the slide
            md_content = md_content_dict.get(slide_number, "")

            # Check if there's no Markdown content for the slide
            if not md_content:
                html.write(f"<div class='content-section'><h2>Slide {slide_number}</h2>")
                html.write("<p class='low-similarity'>No corresponding Markdown file found for this slide.</p>")
                html.write("</div>")
                all_slides_match = False
                continue

            # Calculate similarity
            similarity_score = extract_text_similarity(ppt_content, md_content)

            # Write the slide comparison section
            html.write(f"<div class='content-section'><h2>Slide {slide_number}</h2>")
            if similarity_score <= threshold:
                html.write(f"<p class='low-similarity'>Similarity Score: {similarity_score:.2f}</p>")
                similarity_below_threshold.append({
                    'Slide Number': slide_number,
                    'Similarity Score': similarity_score,
                    'PPT Content': ppt_content,
                    'Markdown Content': md_content
                })
                all_slides_match = False
            else:
                html.write(f"<p class='high-similarity'>Similarity Score: {similarity_score:.2f}</p>")

            # Highlight differences
            ppt_html, md_html = generate_diff_html(ppt_content, md_content)

            html.write("<div class='diff'>")
            html.write("<h3>PowerPoint Content:</h3>")
            html.write(f"<div class='ppt-content'>{ppt_html}</div>")
            html.write("<h3>Markdown Content:</h3>")
            html.write(f"<div class='md-content'>{md_html}</div>")
            html.write("</div></div>")

        # Write the summary section
        html.write("<h2>Summary</h2>")
        if similarity_below_threshold:
            html.write("<p>Slides with similarity below the threshold:</p><ul>")
            for slide in similarity_below_threshold:
                html.write(f"<li>Slide {slide['Slide Number']}: Similarity Score {slide['Similarity Score']:.2f}</li>")
            html.write("</ul>")
        else:
            html.write("<p>All slides have similarity above the threshold.</p>")

        # Write the HTML footer
        html.write("</body></html>")

    print(f"Comparison results have been written to {html_filename}")
    return all_slides_match

def verify_ppt_to_markdown_conversion(ppt_file_path, md_directory_path, threshold=0.75):
    ppt_content_array = extract_ppt_content_with_notes_and_formatting(ppt_file_path)
    md_content_dict = extract_md_content(md_directory_path)

    compare_ppt_to_markdown_advanced(ppt_content_array, md_content_dict, threshold)


def extract_md_content(md_directory_path):
    md_content_dict = {}
    
    for md_file in os.listdir(md_directory_path):
        if md_file.endswith('.md'):
            slide_number = int(md_file.split('_')[0])  # Extract slide number from filename
            md_file_path = os.path.join(md_directory_path, md_file)
            with open(md_file_path, 'r', encoding='utf-8') as file:
                md_content = file.read().strip()
                md_content_dict[slide_number] = md_content
    
    return md_content_dict


def sanitize_text(text, replace):
    """Cleans up unwanted or non-printable characters from the text."""
    # Remove non-printable ASCII characters and excess whitespace
    clean_text = re.sub(
        r"[^\x20-\x7E]", replace, text
    )  # Keeps printable ASCII (space to ~)
    # Removes extra spaces and normalizes
    clean_text = " ".join(clean_text.split())
    return clean_text


def extract_advanced_text_from_slide(slide):
    slide_content = []
    slide_title = ""

    for shape in slide.shapes:
        # Check if the shape contains a text frame
        if shape.has_text_frame:
            if shape == slide.shapes.title:
                slide_title = shape.text.strip()
            else:
                # Handle table shapes separately
                if shape.has_table:
                    table_content = ""
                    # Iterate through each row of the table
                    for row in shape.table.rows:
                        row_text = "| "
                        # Iterate through each cell in the row
                        for cell in row.cells:
                            cell_text = cell.text.strip()  # Strip whitespace from the cell text
                            row_text += f"{cell_text} | "  # Format cell text with pipes
                        table_content += row_text + "\n"  # Append each row to the table content
                    slide_content.append(f"Table:\n{table_content}")  # Append the table content with a label
                
                # Handle special characters or code blocks (use backticks for code)
                elif any(c in shape.text for c in ['`', '{', '}', '(', ')']):
                    # Assume it's a code block or preformatted text
                    slide_content.append(f"``` \n{shape.text.strip()} \n```")
                else:
                    # Handle regular text and paragraphs
                    for paragraph in shape.text_frame.paragraphs:
                        indent_level = paragraph.level
                        text = paragraph.text.strip()
                        formatted_text = f"{'  ' * indent_level}* {text}"
                        slide_content.append(formatted_text)
                    
    return {"title": slide_title, "content": "\n".join(slide_content)}


def extract_ppt_content_with_notes_and_formatting(ppt_file_path):
    ppt_content = []
    presentation = Presentation(ppt_file_path)

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_data = extract_advanced_text_from_slide(slide)
        slide_data['content'] = slide_data['content'].rstrip(f"* {slide_number}")
        
        notes_text = ""
        if slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        # Construct the content for each slide
        content_dict = {
            'Slide_number': slide_number,
            'content': slide_data['title'] + slide_data['content'] + '\nNotes:\n' + notes_text,
        }
        ppt_content.append(content_dict)
    
    return ppt_content

# Example usage
if __name__ == "__main__":
    # Update these paths for your files
    ppt_file_path = 'path/to/your/presentation.pptx'
    md_directory_path = 'path/to/your/markdown/folder'

    result = verify_ppt_to_markdown_conversion(ppt_file_path, md_directory_path)
    if result:
        print('PowerPoint to Markdown conversion is valid!')
    else:
        print('PowerPoint to Markdown conversion has errors!')
