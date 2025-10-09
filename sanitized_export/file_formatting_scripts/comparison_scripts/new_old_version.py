import re
import os
import openpyxl
from pptx import Presentation

def clean_text(text):
    """Remove illegal characters and format title for filenames."""
    text = re.sub(r'[\x00-\x1F\x7F]', '', text) if text else text  # Remove illegal characters
    return text.strip()

def format_filename(text):
    """Format text to be used as a filename: remove spaces, newlines, and special characters."""
    return re.sub(r'[^a-zA-Z0-9]', '', text)  # Keep only alphanumeric characters

def extract_slide_details(prs):
    """Extract slide details: title, text content, and image count."""
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        title = None
        text_content = []
        image_count = 0

        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                try:
                    if shape.placeholder_format.idx == 0:  # Check if it's a title placeholder
                        title = clean_text(shape.text.strip())
                except ValueError:
                    text_content.append(clean_text(shape.text.strip()))  # Non-title text

            if shape.shape_type == 13:  # Shape type 13 corresponds to an image
                image_count += 1

        if title:
            slides.append({"title": title, "slide_number": i, "text": "\n".join(text_content), "images": image_count})

    return slides

def find_closest_match(new_slide_number, old_slides):
    """Find the closest matching old slide within a ±2 slide range."""
    closest_match = None
    min_distance = float("inf")

    for old_slide in old_slides:
        distance = abs(new_slide_number - old_slide["slide_number"])
        if distance <= 10 and distance < min_distance:
            closest_match = old_slide
            min_distance = distance

    return closest_match

def compare_slide_titles_and_content(ppt_old, ppt_new, output_folder):
    """Compare slide titles, text content, and images, ensuring close matches."""
    prs_old = Presentation(ppt_old)
    prs_new = Presentation(ppt_new)

    new_slides = extract_slide_details(prs_new)  # New PPT is the source of truth
    old_slides = extract_slide_details(prs_old)

    # Get first slide title for naming the Excel file
    first_slide_title = new_slides[0]["title"] if new_slides else "ComparisonReport"
    formatted_title = format_filename(first_slide_title)

    # Construct Excel file path
    output_excel = os.path.join(output_folder, f"{formatted_title}.xlsx")

    # Create an Excel workbook
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Slide Comparison"
    ws.append(["Title", "Slide Number (New)", "Slide Number (Old)", "Changes Detected"])

    matched_old_slides = set()

    for new_slide in new_slides:
        title = new_slide["title"]
        new_slide_number = new_slide["slide_number"]

        # Find the closest old slide with the same title
        potential_matches = [slide for slide in old_slides if slide["title"] == title and slide["slide_number"] not in matched_old_slides]
        closest_old_slide = find_closest_match(new_slide_number, potential_matches)

        if closest_old_slide:
            old_slide_number = closest_old_slide["slide_number"]
            matched_old_slides.add(old_slide_number)  # Mark as matched
        else:
            old_slide_number = "Not Found"

        changes = []

        if closest_old_slide:
            # Compare text
            if new_slide["text"] != closest_old_slide["text"]:
                changes.append("Text changed")

            # Compare image count
            if new_slide["images"] != closest_old_slide["images"]:
                changes.append("Image changed")

        change_status = ", ".join(changes) if changes else "No changes"
        ws.append([title, new_slide_number, old_slide_number, change_status])

    # Save the Excel file
    os.makedirs(output_folder, exist_ok=True)
    wb.save(output_excel)
    print(f"Comparison report saved to {output_excel}")

# Example usage:
if __name__ == "__main__":
    # Update these paths for your files
    ppt_old_path = "path/to/your/old/presentation.pptx"
    ppt_new_path = "path/to/your/new/presentation.pptx"
    output_folder = "path/to/output/folder"

    compare_slide_titles_and_content(ppt_old_path, ppt_new_path, output_folder)
