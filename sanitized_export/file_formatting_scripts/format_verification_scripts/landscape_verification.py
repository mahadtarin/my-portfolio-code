import os
import logging
from pptx import Presentation
from langdetect import detect

# Configure logging
log_file_path = "landscape_verification.log"
logging.basicConfig(filename=log_file_path, level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def verify_slide_structure(ppt_path):
    presentation = Presentation(ppt_path)
    slide_number = 0

    for slide in presentation.slides:
        slide_number += 1
        textboxes = []
        images = []

        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                # Exclude title placeholders
                if not (shape.is_placeholder and shape.placeholder_format.idx == 0):
                    textboxes.append(shape)
            elif shape.shape_type == 13:  # Image type
                images.append(shape)

        num_text_columns = len(textboxes)
        issues = []

        # 1. Check column structure
        if num_text_columns not in [1, 2, 3]:
            issues.append(f"Slide {slide_number}: Invalid column structure ({num_text_columns} detected, excluding title).")
        else:
            logging.info(f"Slide {slide_number}: Valid Column Structure.")

        # 2. Validate text constraints
        for i, textbox in enumerate(textboxes):
            text = textbox.text.strip()
            num_lines = text.count("\n") + 1
            language = detect(text)

            # Add any specific text constraints checks here
            # Example: Check if the text is in English
            if language != 'en':
                issues.append(f"Slide {slide_number}, Textbox {i+1}: Non-English text detected.")

        # Log issues
        for issue in issues:
            logging.warning(issue)

if __name__ == "__main__":
    # Update this path for your file
    ppt_path = "path/to/your/presentation.pptx"
    verify_slide_structure(ppt_path)
    print(f"Verification complete. Check the log file at {log_file_path} for details.")